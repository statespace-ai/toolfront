import json
import logging
import xml.etree.ElementTree as ET
from abc import ABC
from pathlib import Path
from typing import Any
from urllib.parse import ParseResult, urlparse

from pydantic import BaseModel, Field, field_validator

from toolfront.models.database import SearchMode
from toolfront.types import ConnectionResult
from toolfront.utils import search_items

logger = logging.getLogger("toolfront")

MAX_SAMPLE_CHARS = 10000
MAX_SAMPLE_ROWS = 50
PAGE_SAMPLE_CHARS = 1000

DOCUMENT_EXTENSIONS = {
    ".docx",
    ".xlsx",
    ".json",
    ".md",
    ".pdf",
    ".pptx",
    ".rtf",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}


class LibraryError(Exception):
    """Exception for library-related errors."""

    pass


class Library(BaseModel, ABC):
    """Abstract base class for library."""

    url: ParseResult = Field(description="URL of the library")

    @field_validator("url", mode="before")
    def validate_url(cls, v: Any) -> ParseResult:
        if isinstance(v, str):
            v = urlparse(v)
        return v

    async def test_connection(self) -> ConnectionResult:
        """Test the connection to the library."""
        return ConnectionResult(connected=True, message="Library connection successful")

    async def get_documents(self) -> list[str]:
        """Get all documents in the library recursively."""
        path = Path(self.url.path)
        if not path.exists():
            return []

        try:
            return [str(p) for p in path.rglob("*.*") if p.suffix.lower() in DOCUMENT_EXTENSIONS]
        except (PermissionError, OSError) as e:
            logger.warning(f"Error accessing {path}: {e}")
            return []

    async def search_documents(self, pattern: str, mode: SearchMode = SearchMode.REGEX, limit: int = 10) -> list[str]:
        """Search for documents in the library."""
        files = await self.get_documents()
        return search_items(files, pattern, mode, limit)

    # DOCX Methods
    def _sample_docx(self, path: str) -> str:
        """Sample DOCX document."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(path)
            text = ""
            for paragraph in doc.paragraphs[:10]:  # Sample first 10 paragraphs
                text += paragraph.text + "\n"
            return text[:MAX_SAMPLE_CHARS]
        except ImportError:
            return "DOCX support requires python-docx library"
        except Exception as e:
            return f"Error reading DOCX file: {str(e)}"

    def _read_docx(self, path: str) -> str:
        """Read DOCX document."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            return "DOCX support requires python-docx library"
        except Exception as e:
            return f"Error reading DOCX file: {str(e)}"

    # Excel Methods
    def _sample_excel(self, path: str) -> str:
        """Sample Excel document."""
        try:
            import pandas as pd

            df = pd.read_excel(path, nrows=MAX_SAMPLE_ROWS)
            return f"Shape: {df.shape}\nColumns: {list(df.columns)}\n\nSample data:\n{df.head().to_string()}"
        except ImportError:
            return "Excel support requires pandas and openpyxl libraries"
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    def _read_excel(self, path: str) -> str:
        """Read Excel document."""
        try:
            import pandas as pd

            excel_data = pd.read_excel(path, sheet_name=None)
            result = ""
            for sheet_name, df in excel_data.items():
                result += f"Sheet: {sheet_name}\n"
                result += f"Shape: {df.shape}\n"
                result += f"Columns: {list(df.columns)}\n"
                result += df.to_string() + "\n\n"
            return result
        except ImportError:
            return "Excel support requires pandas and openpyxl libraries"
        except Exception as e:
            return f"Error reading Excel file: {str(e)}"

    # JSON Methods
    def _sample_json(self, path: str) -> str:
        """Sample JSON document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                data = json.load(file)
                return json.dumps(data, indent=2, ensure_ascii=False)[:MAX_SAMPLE_CHARS]
        except json.JSONDecodeError as e:
            return f"Error parsing JSON: {str(e)}"
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"

    def _read_json(self, path: str) -> str:
        """Read JSON document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                data = json.load(file)
                return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError as e:
            return f"Error parsing JSON: {str(e)}"
        except Exception as e:
            return f"Error reading JSON file: {str(e)}"

    # Markdown Methods
    def _sample_markdown(self, path: str) -> str:
        """Sample Markdown document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                text = file.read()
                return text[:MAX_SAMPLE_CHARS]
        except Exception as e:
            return f"Error reading Markdown file: {str(e)}"

    def _read_markdown(self, path: str) -> str:
        """Read Markdown document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            return f"Error reading Markdown file: {str(e)}"

    # PDF Methods
    def _sample_pdf(self, path: str) -> str:
        """Sample PDF document."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(path)
            text = ""
            total_chars = 0

            for page_num, page in enumerate(reader.pages, 1):
                # Sample first PAGE_SAMPLE_CHARS chars of page
                page_text = page.extract_text()[:PAGE_SAMPLE_CHARS]
                text += f"Page {page_num}: {page_text}...\n"
                total_chars += len(text)

                if total_chars >= MAX_SAMPLE_CHARS:
                    text = text[:MAX_SAMPLE_CHARS]
                    break

            return text
        except ImportError:
            return "PDF support requires pypdf library"
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    def _read_pdf(self, path: str) -> str:
        """Read PDF document."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        except ImportError:
            return "PDF support requires pypdf library"
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    # PowerPoint Methods
    def _sample_powerpoint(self, path: str) -> str:
        """Sample PowerPoint document."""
        try:
            from pptx import Presentation

            prs = Presentation(path)
            text = ""
            total_chars = 0

            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                slide_text += "\n"

                # Sample first PAGE_SAMPLE_CHARS chars from slide
                slide_sample = slide_text[:PAGE_SAMPLE_CHARS]
                text += slide_sample
                total_chars += len(slide_sample)

                if total_chars >= MAX_SAMPLE_CHARS:
                    text = text[:MAX_SAMPLE_CHARS]
                    break

            return text
        except ImportError:
            return "PowerPoint support requires python-pptx library"
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"

    def _read_powerpoint(self, path: str) -> str:
        """Read PowerPoint document."""
        try:
            from pptx import Presentation

            prs = Presentation(path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"Slide {i + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except ImportError:
            return "PowerPoint support requires python-pptx library"
        except Exception as e:
            return f"Error reading PowerPoint file: {str(e)}"

    # RTF Methods
    def _sample_rtf(self, path: str) -> str:
        """Sample RTF document."""
        try:
            from striprtf.striprtf import rtf_to_text

            with Path(path).open("r", encoding="utf-8") as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return text[:MAX_SAMPLE_CHARS]
        except ImportError:
            # Fallback to raw RTF if striprtf not available
            with Path(path).open("r", encoding="utf-8") as file:
                content = file.read()[:MAX_SAMPLE_CHARS]
                return f"RTF content (requires striprtf library for text extraction):\n{content}"
        except Exception as e:
            return f"Error reading RTF file: {str(e)}"

    def _read_rtf(self, path: str) -> str:
        """Read RTF document."""
        try:
            from striprtf.striprtf import rtf_to_text

            with Path(path).open("r", encoding="utf-8") as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return text
        except ImportError:
            # Fallback to raw RTF if striprtf not available
            with Path(path).open("r", encoding="utf-8") as file:
                return f"RTF content (requires striprtf library for text extraction):\n{file.read()}"
        except Exception as e:
            return f"Error reading RTF file: {str(e)}"

    # TXT Methods
    def _sample_txt(self, path: str) -> str:
        """Sample TXT document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()[:MAX_SAMPLE_CHARS]
        except Exception as e:
            return f"Error reading TXT file: {str(e)}"

    def _read_txt(self, path: str) -> str:
        """Read TXT document."""
        try:
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            return f"Error reading TXT file: {str(e)}"

    # XML Methods
    def _sample_xml(self, path: str) -> str:
        """Sample XML document."""
        try:
            tree = ET.parse(path)
            root = tree.getroot()

            result = f"Root element: {root.tag}\n"
            result += f"Root attributes: {root.attrib}\n\n"

            for i, child in enumerate(root):
                if i >= 5:  # Limit to first 5 elements
                    result += "... (truncated)\n"
                    break
                result += f"Element: {child.tag}\n"
                result += f"Attributes: {child.attrib}\n"
                if child.text and child.text.strip():
                    result += f"Text: {child.text.strip()}\n"
                result += "\n"

            return result[:MAX_SAMPLE_CHARS]
        except Exception:
            # Fallback to plain text if parsing fails
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()[:MAX_SAMPLE_CHARS]

    def _read_xml(self, path: str) -> str:
        """Read XML document."""
        try:
            tree = ET.parse(path)
            root = tree.getroot()

            def element_to_string(elem, level=0):
                indent = "  " * level
                result = f"{indent}Element: {elem.tag}\n"
                if elem.attrib:
                    result += f"{indent}Attributes: {elem.attrib}\n"
                if elem.text and elem.text.strip():
                    result += f"{indent}Text: {elem.text.strip()}\n"
                for child in elem:
                    result += element_to_string(child, level + 1)
                return result

            return element_to_string(root)
        except Exception:
            # Fallback to plain text if parsing fails
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()

    # YAML Methods
    def _sample_yaml(self, path: str) -> str:
        """Sample YAML document."""
        try:
            import yaml

            with Path(path).open("r", encoding="utf-8") as file:
                data = yaml.safe_load(file)
                return yaml.dump(data, indent=2, default_flow_style=False)[:MAX_SAMPLE_CHARS]
        except ImportError:
            # Fallback to plain text if PyYAML not available
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()[:MAX_SAMPLE_CHARS]
        except Exception as e:
            return f"Error parsing YAML: {str(e)}"

    def _read_yaml(self, path: str) -> str:
        """Read YAML document."""
        try:
            import yaml

            with Path(path).open("r", encoding="utf-8") as file:
                data = yaml.safe_load(file)
                return yaml.dump(data, indent=2, default_flow_style=False)
        except ImportError:
            # Fallback to plain text if PyYAML not available
            with Path(path).open("r", encoding="utf-8") as file:
                return file.read()
        except Exception as e:
            return f"Error parsing YAML: {str(e)}"

    async def sample_document(self, path: str) -> str:
        """Sample the file using appropriate method based on file extension."""
        file_extension = path.split(".")[-1].lower()

        # Switch based on file extension
        if file_extension == "docx":
            return self._sample_docx(path)
        elif file_extension == "xlsx":
            return self._sample_excel(path)
        elif file_extension == "json":
            return self._sample_json(path)
        elif file_extension == "md":
            return self._sample_markdown(path)
        elif file_extension == "pdf":
            return self._sample_pdf(path)
        elif file_extension == "pptx":
            return self._sample_powerpoint(path)
        elif file_extension == "rtf":
            return self._sample_rtf(path)
        elif file_extension == "txt":
            return self._sample_txt(path)
        elif file_extension == "xml":
            return self._sample_xml(path)
        elif file_extension == "yaml" or file_extension == "yml":
            return self._sample_yaml(path)
        else:
            return f"Unsupported file type: {file_extension}"

    async def read_document(self, path: str) -> str:
        """Read the file using appropriate method based on file extension."""
        file_extension = path.split(".")[-1].lower()

        # Switch based on file extension
        if file_extension == "docx":
            return self._read_docx(path)
        elif file_extension == "xlsx":
            return self._read_excel(path)
        elif file_extension == "json":
            return self._read_json(path)
        elif file_extension == "md":
            return self._read_markdown(path)
        elif file_extension == "pdf":
            return self._read_pdf(path)
        elif file_extension == "pptx":
            return self._read_powerpoint(path)
        elif file_extension == "rtf":
            return self._read_rtf(path)
        elif file_extension == "txt":
            return self._read_txt(path)
        elif file_extension == "xml":
            return self._read_xml(path)
        elif file_extension == "yaml" or file_extension == "yml":
            return self._read_yaml(path)
        else:
            return f"Unsupported file type: {file_extension}"
