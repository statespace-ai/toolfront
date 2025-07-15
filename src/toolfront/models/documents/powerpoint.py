import logging

from toolfront.models.documents.base import Document, DocumentError

logger = logging.getLogger("toolfront")


class PowerPointDocument(Document):
    """PowerPoint document reader."""

    async def read(self, pagination: int | float = 0) -> str:
        """Read PowerPoint document content.

        Args:
            pagination: Slide number (1+ int) or percentile (0-1 exclusive float) to read.

        Returns:
            Document content as string.
        """
        try:
            from pptx import Presentation

            prs = Presentation(self.path)
            total_slides = len(prs.slides)

            target_slide_idx = self._get_target_page(pagination, total_slides) - 1

            slide = prs.slides[target_slide_idx]
            text = f"Slide {target_slide_idx + 1} of {total_slides}:\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

            return text
        except ImportError:
            error_msg = "PowerPoint support requires python-pptx library"
            logger.error(f"Missing dependency for PowerPoint reading: {error_msg}")
            raise DocumentError(error_msg)
        except Exception as e:
            logger.error(f"Error reading PowerPoint file {self.url}: {e}")
            raise DocumentError(f"Error reading PowerPoint file: {str(e)}")
